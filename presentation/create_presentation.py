from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from datetime import date

# Create a new presentation
prs = Presentation()

# Helper function to add a slide with title and content
def add_slide(title_text, content_text, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]  # using a Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    # Format content text with bullet points
    body = slide.placeholders[1].text_frame
    body.clear()  # clear any existing paragraphs
    for line in content_text.split("\n"):
        p = body.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
        p.level = 0
    return slide

# --------- Slide 1: Title Slide ---------
title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Crypto Bank Project"
# Subtitle placeholder (typically placeholder 1 in the title slide)
slide.placeholders[1].text = "Your Name\nDate: {:%B %d, %Y}\nCourse: Advanced Web Projects".format(date.today())

# --------- Slide 2: Introduction ---------
intro_text = (
    "Overview:\n"
    "• A secure and efficient crypto banking solution\n"
    "• Built using Node.js for the backend and Next.js for the frontend\n"
    "• Designed to support both crypto and fiat transactions seamlessly"
)
add_slide("Introduction", intro_text)

# --------- Slide 3: Features Overview ---------
features_text = (
    "Key Features:\n"
    "• User Account Management\n"
    "• Integrated Wallet System\n"
    "• KYC Verification Process\n"
    "• Merchant Payment Processing\n"
    "• Admin & Staff Management Tools\n"
    "• Transaction Processing (Crypto & Fiat)"
)
add_slide("Features Overview", features_text)

# --------- Slide 4: Class Diagram ---------
diagram_text = (
    "Class Diagram Details:\n"
    "• Main entities: User, Wallet, KYC, Merchant, Admin, Staff, Transaction\n"
    "• Relationships & key methods illustrated\n"
    "• (Insert diagram image in the slide for visual reference)"
)
slide_diagram = add_slide("Class Diagram", diagram_text)
# Optional: Uncomment the following lines and update the path to insert your diagram image
# left = Inches(1)
# top = Inches(3.5)
# width = Inches(8)
# slide_diagram.shapes.add_picture("class_diagram.png", left, top, width=width)

# --------- Slide 5: System Architecture ---------
architecture_text = (
    "Architecture Overview:\n"
    "• Frontend: Next.js for dynamic user interfaces\n"
    "• Backend: Node.js with RESTful API endpoints\n"
    "• Database: SQL/NoSQL options depending on needs\n"
    "• Security: JWT authentication and data encryption\n"
    "• Communication: API-driven interactions between services"
)
add_slide("System Architecture", architecture_text)

# --------- Slide 6: User Flow ---------
user_flow_text = (
    "User Journey:\n"
    "1. User Registration & KYC Verification\n"
    "2. Wallet Creation & Funding\n"
    "3. Initiation of Crypto/Fiat Transactions\n"
    "4. Merchant Payment Processing\n"
    "5. Ongoing Admin/Staff Management"
)
add_slide("User Flow", user_flow_text)

# --------- Slide 7: Technologies Used ---------
technologies_text = (
    "Technology Stack:\n"
    "• Frontend: Next.js, Tailwind CSS (or another UI framework)\n"
    "• Backend: Node.js (Express or NestJS)\n"
    "• Database: MongoDB / PostgreSQL\n"
    "• Authentication: JWT, OAuth\n"
    "• Payment Integration: APIs (e.g., Binance, Coinbase)\n"
    "• DevOps: Docker, CI/CD pipelines"
)
add_slide("Technologies Used", technologies_text)

# --------- Slide 8: Challenges & Solutions ---------
challenges_text = (
    "Project Challenges:\n"
    "• Ensuring robust security and data encryption\n"
    "• Optimizing transaction speed and reliability\n"
    "• Scalability to handle growing user base\n\n"
    "Solutions:\n"
    "• Implementing state-of-the-art security protocols\n"
    "• Performance tuning of API endpoints\n"
    "• Designing with scalability in mind (microservices, load balancing)"
)
add_slide("Challenges & Solutions", challenges_text)

# --------- Slide 9: Future Enhancements ---------
enhancements_text = (
    "Planned Enhancements:\n"
    "• Multi-chain support for broader crypto integration\n"
    "• Additional financial features (staking, lending)\n"
    "• Advanced analytics and reporting tools\n"
    "• Continuous UI/UX improvements based on user feedback"
)
add_slide("Future Enhancements", enhancements_text)

# --------- Slide 10: Conclusion ---------
conclusion_text = (
    "Conclusion:\n"
    "• The Crypto Bank project aims to revolutionize digital transactions.\n"
    "• Focus on security, scalability, and user experience sets it apart.\n"
    "• A solid foundation built with modern technologies paves the way for future innovation."
)
add_slide("Conclusion", conclusion_text)

# --------- Slide 11: Demo (Optional) ---------
demo_text = (
    "Demo Overview:\n"
    "• Screenshots or a short video demo of the application\n"
    "• Key interactions and user interface highlights\n"
    "• Live demo URL if available"
)
add_slide("Demo", demo_text)

# Save the presentation to a file
output_filename = "CryptoBankPresentation.pptx"
prs.save(output_filename)
print(f"Presentation created successfully as '{output_filename}'")